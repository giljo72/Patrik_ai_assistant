import os
import pdfplumber
import docx
from pptx import Presentation

def load_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    elif ext == ".pdf":
        with pdfplumber.open(file_path) as pdf:
            return "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())

    elif ext == ".docx":
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    else:
        raise ValueError(f"Unsupported file type: {ext}")
